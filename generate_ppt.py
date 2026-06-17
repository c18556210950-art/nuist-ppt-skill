"""
南京信息工程大学 PPT 生成器 - 简单模式
直接在模板上复制版式、填写内容。

用法：
    1. 修改下方「配置区」的内容
    2. 运行: python generate_ppt.py

依赖：
    pip install python-pptx
"""

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
import os

# ═══════════════════════════════════════════════════════
# 配置区：修改这里的内容即可
# ═══════════════════════════════════════════════════════

OUTPUT_FILE = "output.pptx"

COVER = {
    "title": "毕业论文答辩",
    "subtitle": "南京信息工程大学",
    "student": "张三",
    "advisor": "李四 教授",
}

TOC_ITEMS = [
    "研究背景与意义",
    "相关研究工作",
    "研究方法与技术路线",
    "实验结果与分析",
    "总结与展望",
]

SECTIONS = [
    {"num": "1", "title": "研究背景与意义", "subtitle": "Background & Significance"},
    {"num": "2", "title": "研究方法与技术路线", "subtitle": "Methodology & Approach"},
    {"num": "3", "title": "实验结果与分析", "subtitle": "Results & Analysis"},
    {"num": "4", "title": "总结与展望", "subtitle": "Summary & Outlook"},
]

CONTENT_SLIDES = [
    {
        "title": "研究背景",
        "bullets": [
            "云计算技术快速发展，边缘计算成为研究热点",
            "物联网设备数量指数级增长，对计算资源需求日益增加",
            "传统云计算架构存在延迟高、带宽消耗大等问题",
        ],
    },
    {
        "title": "国内外研究现状",
        "bullets": [
            "国外研究：Google Edge TPU、AWS Wavelength 等商业方案已落地",
            "国内研究：华为云边缘计算平台、阿里云 Link Edge 等快速发展",
            "学术界关注：资源调度、服务迁移、安全隐私等关键问题",
        ],
    },
    {
        "title": "技术路线",
        "bullets": [
            "数据采集：传感器数据、网络状态、用户行为",
            "模型构建：深度学习 + 强化学习 + 优化模型",
            "实验验证：仿真平台 + 真实场景 + 对比分析",
        ],
    },
    {
        "title": "主要创新点",
        "bullets": [
            "提出了一种基于深度强化学习的边缘计算资源调度算法",
            "设计了多目标优化的服务迁移策略",
            "在真实边缘计算平台上进行了大规模实验验证",
        ],
    },
]

CLOSING = {
    "thanks": "感谢各位老师指导！",
    "subtitle": "Thank you for your support!",
    "student": "张三",
    "advisor": "李四 教授",
}

TEMPLATE_FILE = os.path.join(os.path.dirname(__file__), "南京信息工程大学ppt模版.pptx")

# ═══════════════════════════════════════════════════════
# 生成逻辑
# ═══════════════════════════════════════════════════════

TITLE = 1
BODY = 2
CENTER_TITLE = 3
SUBTITLE = 4
OBJECT = 7


def find_layout(prs, keyword):
    for layout in prs.slide_layouts:
        if keyword in layout.name:
            return layout
    return None


def fill_ph(slide, ph_type, text, idx=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type:
            if idx is not None and ph.placeholder_format.idx != idx:
                continue
            ph.text = text
            return ph
    return None


def get_body_ph(slide):
    candidates = [ph for ph in slide.placeholders
                  if ph.placeholder_format.type in (BODY, OBJECT)]
    return max(candidates, key=lambda ph: ph.width) if candidates else None


def set_bullets(ph, texts):
    tf = ph.text_frame
    tf.clear()
    for i, text in enumerate(texts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0


def remove_template_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    while len(sldIdLst) > 0:
        rId = sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        del sldIdLst[0]


def add_cover(prs):
    slide = prs.slides.add_slide(find_layout(prs, "标题幻灯片"))
    ph = fill_ph(slide, CENTER_TITLE, COVER["title"])
    if ph: ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ph = fill_ph(slide, SUBTITLE, COVER["subtitle"])
    if ph: ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for ph in slide.placeholders:
        if ph.placeholder_format.type == BODY:
            if ph.placeholder_format.idx == 17:
                ph.text = f">> 答辩学生：{COVER['student']}"
            elif ph.placeholder_format.idx == 18:
                ph.text = f">> 指导教师：{COVER['advisor']}"
    return slide


def add_toc(prs):
    slide = prs.slides.add_slide(find_layout(prs, "目录1"))
    ph = get_body_ph(slide)
    if ph:
        ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        set_bullets(ph, TOC_ITEMS)
    return slide


def add_section(prs, num, title, subtitle):
    slide = prs.slides.add_slide(find_layout(prs, "节标题"))
    ph = fill_ph(slide, TITLE, title)
    if ph: ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for ph in slide.placeholders:
        if ph.placeholder_format.type == BODY:
            if ph.placeholder_format.idx == 13:
                ph.text = num
                if len(num) >= 2:
                    orig_top = ph.top
                    orig_height = ph.height
                    ph.width = int(ph.width * 2.0)
                    ph.left -= 1651000
                    ph.top = orig_top
                    ph.height = orig_height
                    for p in ph.text_frame.paragraphs:
                        p.alignment = 2  # PP_ALIGN.RIGHT
            elif ph.placeholder_format.idx == 1:
                ph.text = subtitle
                ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return slide


def add_content(prs, title, bullets):
    slide = prs.slides.add_slide(find_layout(prs, "标题和内容"))
    fill_ph(slide, TITLE, title)
    ph = get_body_ph(slide)
    if ph:
        set_bullets(ph, bullets)
    return slide


def add_closing(prs):
    slide = prs.slides.add_slide(find_layout(prs, "结尾幻灯片"))
    fill_ph(slide, CENTER_TITLE, CLOSING["thanks"])
    fill_ph(slide, SUBTITLE, CLOSING["subtitle"])
    for ph in slide.placeholders:
        if ph.placeholder_format.type == BODY:
            if ph.placeholder_format.idx == 17:
                ph.text = f">> 答辩学生：{CLOSING['student']}"
            elif ph.placeholder_format.idx == 18:
                ph.text = f">> 指导教师：{CLOSING['advisor']}"
    return slide


def generate():
    print(f"[打开] {TEMPLATE_FILE}")
    prs = Presentation(TEMPLATE_FILE)
    remove_template_slides(prs)

    print("[1] 封面")
    add_cover(prs)

    print("[2] 目录")
    add_toc(prs)

    for sec in SECTIONS:
        print(f"[3] 章节: {sec['title']}")
        add_section(prs, sec["num"], sec["title"], sec["subtitle"])

    for cs in CONTENT_SLIDES:
        print(f"[4] 内容: {cs['title']}")
        add_content(prs, cs["title"], cs["bullets"])

    print("[5] 结尾")
    add_closing(prs)

    print(f"[保存] {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print(f"[完成] {len(prs.slides)} slides -> {OUTPUT_FILE}")


if __name__ == "__main__":
    generate()
